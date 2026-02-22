"""Synapse Verification Engine — API Routes

All /api/* endpoints for the claim-level intelligence pipeline.
Extracted from main.py to keep the core product endpoints isolated
from legacy hackathon features.
"""
from __future__ import annotations

import os
import io
import uuid
import datetime
from typing import List, Optional, Dict, Any

from fastapi import APIRouter, HTTPException, UploadFile, File
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from app.verification_engine import (
    extract_claims, extract_url_content, run_verification_pipeline,
    trace_provenance, generate_corrected_claim, VerificationEvent,
)

try:
    from deepgram import DeepgramClient, PrerecordedOptions  # type: ignore
except Exception:
    DeepgramClient = None  # type: ignore
    PrerecordedOptions = None  # type: ignore


router = APIRouter(prefix="/api", tags=["synapse"])

# ---------------------------------------------------------------------------
# In-memory report store (hackathon — no DB needed)
# ---------------------------------------------------------------------------
_reports_store: Dict[str, dict] = {}


# ---------------------------------------------------------------------------
# Request / Response Models
# ---------------------------------------------------------------------------

class IngestRequest(BaseModel):
    url: Optional[str] = None
    text: Optional[str] = None

class IngestResponse(BaseModel):
    title: str
    text: str
    source_type: str  # url, text, audio, sec_filing, earnings_transcript, financial_document
    url: Optional[str] = None

class ExtractClaimsRequest(BaseModel):
    text: str

class ClaimLocation(BaseModel):
    chunk_id: str
    start_char: int
    end_char: int

class ClaimItem(BaseModel):
    id: str
    original: str
    normalized: str
    type: str
    location: Optional[ClaimLocation] = None
    location_str: str = ""  # backward-compat "c0001:120-185"
    company_ticker: Optional[str] = None

class ExtractClaimsResponse(BaseModel):
    claims: List[ClaimItem]

class VerifyRequest(BaseModel):
    claim: str

class SaveReportRequest(BaseModel):
    title: str
    url: Optional[str] = None
    source_type: str = "url"
    claims: list  # full claim objects with verification data
    analyzed_at: Optional[str] = None

class AuditLogRequest(BaseModel):
    title: str
    url: Optional[str] = None
    claims: list
    analyzed_at: Optional[str] = None


# ---------------------------------------------------------------------------
# Reports
# ---------------------------------------------------------------------------

@router.post("/reports")
def api_save_report(req: SaveReportRequest):
    """Save a verification report and return a unique shareable ID."""
    report_id = str(uuid.uuid4())[:8]
    _reports_store[report_id] = {
        "id": report_id,
        "title": req.title,
        "url": req.url,
        "source_type": req.source_type,
        "claims": req.claims,
        "analyzed_at": req.analyzed_at or "",
        "created_at": datetime.datetime.utcnow().isoformat(),
    }
    return {"id": report_id}


@router.get("/reports/{report_id}")
def api_get_report(report_id: str):
    """Retrieve a saved verification report by ID."""
    report = _reports_store.get(report_id)
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    return report


# ---------------------------------------------------------------------------
# Audit Log Export
# ---------------------------------------------------------------------------

@router.post("/export-audit-log")
def api_export_audit_log(req: AuditLogRequest):
    """Generate a structured audit trail for compliance purposes."""
    audit_log = {
        "compliance_metadata": {
            "verification_id": str(uuid.uuid4()),
            "timestamp": datetime.datetime.utcnow().isoformat(),
            "engine_version": "synapse-financial-1.0",
            "source_title": req.title,
            "source_url": req.url,
            "analyzed_at": req.analyzed_at,
        },
        "claims": [],
    }
    for claim in req.claims:
        claim_entry = {
            "claim_text": claim.get("original", ""),
            "claim_type": claim.get("type", ""),
            "verdict": None,
            "confidence": None,
            "verified_against": None,
            "evidence_chain": [],
            "contradictions": [],
        }
        v = claim.get("verification", {})
        if v:
            ov = v.get("overallVerdict", {})
            claim_entry["verdict"] = ov.get("verdict")
            claim_entry["confidence"] = ov.get("confidence")
            claim_entry["verified_against"] = ov.get("verified_against")
            for ev in v.get("evidence", []):
                claim_entry["evidence_chain"].append({
                    "source_id": ev.get("id"),
                    "title": ev.get("title"),
                    "tier": ev.get("tier"),
                    "source_url": ev.get("source"),
                    "filing_type": ev.get("filing_type"),
                    "accession_number": ev.get("accession_number"),
                    "filing_date": ev.get("filing_date"),
                    "quality_score": ev.get("quality_score"),
                    "supports_claim": ev.get("supports_claim"),
                })
            for c in v.get("contradictions", []):
                claim_entry["contradictions"].append(c)
        audit_log["claims"].append(claim_entry)
    return audit_log


# ---------------------------------------------------------------------------
# Ingest (URL, text, audio, file)
# ---------------------------------------------------------------------------

@router.post("/ingest", response_model=IngestResponse)
def api_ingest(req: IngestRequest):
    """Ingest content from URL or raw text. Returns clean extracted text."""
    if req.url:
        # Detect SEC filing URLs
        source_type = "url"
        if "sec.gov" in req.url.lower():
            source_type = "sec_filing"
        elif any(kw in req.url.lower() for kw in ["earnings", "transcript"]):
            source_type = "earnings_transcript"

        result = extract_url_content(req.url)
        if result.get("error"):
            raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {result['error']}")
        return IngestResponse(
            title=result.get("title", req.url),
            text=result.get("text", ""),
            source_type=source_type,
            url=req.url,
        )
    elif req.text:
        # Direct text paste — just pass through
        title = req.text[:60].strip().replace("\n", " ")
        return IngestResponse(title=title, text=req.text, source_type="text")
    else:
        raise HTTPException(status_code=400, detail="Provide either 'url' or 'text'")


@router.post("/ingest-audio", response_model=IngestResponse)
async def api_ingest_audio(file: UploadFile = File(...)):
    """Ingest audio/video file — transcribe via Deepgram, return text."""
    deepgram_key = os.getenv("DEEPGRAM_API_KEY")
    if not deepgram_key:
        raise HTTPException(status_code=503, detail="Deepgram API not configured")

    audio_bytes = await file.read()
    try:
        client = DeepgramClient(deepgram_key)
        options = PrerecordedOptions(
            model="nova-2", language="en", smart_format=True,
            punctuate=True, paragraphs=True, diarize=True,
        )
        source = {"buffer": audio_bytes, "mimetype": file.content_type or "audio/mpeg"}
        response = client.listen.rest.v("1").transcribe_file(source, options)
        transcript = response.results.channels[0].alternatives[0].transcript
        return IngestResponse(
            title=file.filename or "Audio Upload",
            text=transcript,
            source_type="audio",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Transcription failed: {e}")


@router.post("/ingest-file", response_model=IngestResponse)
async def api_ingest_file(file: UploadFile = File(...)):
    """Ingest a document file (PDF, PPTX, DOCX) — extract text and return it."""
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")

    ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
    file_bytes = await file.read()

    if ext == "pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            pages = []
            for page in doc:
                pages.append(page.get_text())
            doc.close()
            text = "\n\n".join(pages).strip()
            if not text:
                raise HTTPException(status_code=400, detail="PDF appears to contain no extractable text (may be scanned/image-only)")
            title = file.filename
            # Try to get title from PDF metadata
            meta = fitz.open(stream=file_bytes, filetype="pdf").metadata
            if meta and meta.get("title"):
                title = meta["title"]
            return IngestResponse(title=title, text=text, source_type="pdf")
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PDF parsing failed: {e}")

    elif ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_parts = [f"--- Slide {i} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_parts.append(t)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip(" |"):
                                slide_parts.append(row_text)
                # Also grab speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"[Speaker Notes] {notes}")
                slides_text.append("\n".join(slide_parts))
            text = "\n\n".join(slides_text).strip()
            if not text or text.replace("-", "").replace("Slide", "").strip() == "":
                raise HTTPException(status_code=400, detail="PowerPoint appears to contain no extractable text")
            title = file.filename
            if prs.core_properties and prs.core_properties.title:
                title = prs.core_properties.title
            return IngestResponse(title=title, text=text, source_type="pptx")
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPTX parsing failed: {e}")

    elif ext in ("docx", "doc"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = []
            for para in doc.paragraphs:
                t = para.text.strip()
                if t:
                    paragraphs.append(t)
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        paragraphs.append(row_text)
            text = "\n\n".join(paragraphs).strip()
            if not text:
                raise HTTPException(status_code=400, detail="Word document appears to contain no extractable text")
            title = file.filename
            if doc.core_properties and doc.core_properties.title:
                title = doc.core_properties.title
            return IngestResponse(title=title, text=text, source_type="docx")
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"DOCX parsing failed: {e}")

    else:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: .{ext}. Supported: .pdf, .pptx, .docx"
        )


# ---------------------------------------------------------------------------
# Extract Claims
# ---------------------------------------------------------------------------

@router.post("/extract-claims", response_model=ExtractClaimsResponse)
def api_extract_claims(req: ExtractClaimsRequest):
    """Extract verifiable factual claims from text."""
    if not req.text.strip():
        raise HTTPException(status_code=400, detail="Text is empty")
    raw_claims = extract_claims(req.text)
    claims = []
    for i, c in enumerate(raw_claims):
        loc = c.get("location")
        claim_location = None
        if isinstance(loc, dict) and "chunk_id" in loc:
            claim_location = ClaimLocation(
                chunk_id=loc["chunk_id"],
                start_char=loc.get("start_char", 0),
                end_char=loc.get("end_char", 0),
            )
        claims.append(ClaimItem(
            id=c.get("id", f"claim-{i+1}"),
            original=c.get("original", ""),
            normalized=c.get("normalized", c.get("original", "")),
            type=c.get("type", "categorical"),
            location=claim_location,
            location_str=c.get("location_str", ""),
            company_ticker=c.get("company_ticker"),
        ))
    return ExtractClaimsResponse(claims=claims)


# ---------------------------------------------------------------------------
# Verify (SSE Stream)
# ---------------------------------------------------------------------------

@router.post("/verify")
def api_verify(req: VerifyRequest):
    """Run the full multi-stage verification pipeline on a single claim.
    Returns Server-Sent Events (SSE) stream for real-time UI updates."""
    if not req.claim.strip():
        raise HTTPException(status_code=400, detail="Claim is empty")

    def event_stream():
        for event in run_verification_pipeline(req.claim):
            yield event.to_sse()

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


# ---------------------------------------------------------------------------
# Financial Claims Feed (demo data)
# ---------------------------------------------------------------------------

@router.get("/financial-claims-feed")
def api_financial_claims_feed():
    """Return hardcoded demo financial claims showcasing verification scenarios."""
    return {"claims": [
        {"id": "fc-1", "text": "Apple's gross margin expanded to 46.2% in Q4 FY2024, up from 45.2% a year ago", "company": "AAPL", "type": "financial_metric", "source": "10-K FY2024"},
        {"id": "fc-2", "text": "Microsoft completed the $68.7B Activision Blizzard acquisition in Oct 2023", "company": "MSFT", "type": "transaction", "source": "8-K Filing"},
        {"id": "fc-3", "text": "Nvidia's data center revenue reached $47.5B in FY2024, up 217% YoY", "company": "NVDA", "type": "financial_metric", "source": "10-K FY2024"},
        {"id": "fc-4", "text": "JPMorgan's CET1 capital ratio was 15.0% as of Q4 2024", "company": "JPM", "type": "regulatory", "source": "10-Q"},
        {"id": "fc-5", "text": "Tesla delivered 1.81 million vehicles globally in 2023, missing its 2M target", "company": "TSLA", "type": "financial_metric", "source": "Earnings Call"},
        {"id": "fc-6", "text": "Goldman Sachs' trading revenue surged to $6.6B in Q3 2024", "company": "GS", "type": "financial_metric", "source": "10-Q"},
        {"id": "fc-7", "text": "Amazon Web Services generated $90.8B in revenue in 2023", "company": "AMZN", "type": "financial_metric", "source": "10-K FY2023"},
        {"id": "fc-8", "text": "Broadcom completed its $61B acquisition of VMware in Nov 2023", "company": "AVGO", "type": "transaction", "source": "8-K Filing"},
    ]}


@router.get("/trending-tweets")
def api_trending_tweets():
    """Redirect to financial claims feed."""
    return api_financial_claims_feed()
